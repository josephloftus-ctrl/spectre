"""
Corpus ingestion for the knowledge RAG database.
Handles PDF, DOCX, PPTX, XLSX files from the Training folder.
Separate from inventory processing pipeline.
"""
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime
import json

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

from .embeddings import embed_text, chunk_text, get_collection

logger = logging.getLogger(__name__)

# Training corpus location
TRAINING_DIR = Path(__file__).parent.parent.parent / "Training"
KNOWLEDGE_COLLECTION = "culinart_bible"


def parse_pdf(file_path: Path) -> str:
    """Extract text from PDF."""
    text_parts = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text.strip())
    except Exception as e:
        logger.error(f"Error parsing PDF {file_path}: {e}")
    return "\n\n".join(text_parts)


def parse_docx(file_path: Path) -> str:
    """Extract text from DOCX."""
    text_parts = []
    try:
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        # Also get text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
    except Exception as e:
        logger.error(f"Error parsing DOCX {file_path}: {e}")
    return "\n\n".join(text_parts)


def parse_pptx(file_path: Path) -> str:
    """Extract text from PPTX."""
    text_parts = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    except Exception as e:
        logger.error(f"Error parsing PPTX {file_path}: {e}")
    return "\n\n".join(text_parts)


def parse_xlsx(file_path: Path) -> str:
    """Extract text from XLSX."""
    text_parts = []
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[Sheet: {sheet_name}]"]
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    sheet_text.append(row_text)
            if len(sheet_text) > 1:
                text_parts.append("\n".join(sheet_text))
        wb.close()
    except Exception as e:
        logger.error(f"Error parsing XLSX {file_path}: {e}")
    return "\n\n".join(text_parts)


def parse_file(file_path: Path) -> Optional[str]:
    """Parse any supported file type and return text content."""
    suffix = file_path.suffix.lower()

    if suffix == '.pdf':
        return parse_pdf(file_path)
    elif suffix == '.docx':
        return parse_docx(file_path)
    elif suffix == '.pptx':
        return parse_pptx(file_path)
    elif suffix == '.xlsx':
        return parse_xlsx(file_path)
    else:
        logger.warning(f"Unsupported file type: {suffix} for {file_path}")
        return None


def ingest_file(file_path: Path, collection_name: str = KNOWLEDGE_COLLECTION) -> Dict[str, Any]:
    """
    Ingest a single file into the knowledge RAG database.

    Returns:
        Dict with success status and chunk count
    """
    collection = get_collection(collection_name)
    if not collection:
        return {"success": False, "error": "Collection not available"}

    # Parse the file
    content = parse_file(file_path)
    if not content:
        return {"success": False, "error": "Could not parse file"}

    # Create chunks
    chunks = chunk_text(content, max_size=500, overlap=50)

    # Generate embeddings and store
    embedded_count = 0
    file_id = file_path.stem  # Use filename without extension as ID

    for i, chunk in enumerate(chunks):
        chunk_id = f"knowledge_{file_id}_{i}"

        embedding = embed_text(chunk)
        if not embedding:
            continue

        try:
            collection.upsert(
                ids=[chunk_id],
                embeddings=[embedding],
                documents=[chunk],
                metadatas=[{
                    "source_file": file_path.name,
                    "source_type": "knowledge_corpus",
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "ingested_at": datetime.now().isoformat(),
                    "collection": collection_name
                }]
            )
            embedded_count += 1
        except Exception as e:
            logger.error(f"Failed to store chunk {chunk_id}: {e}")

    return {
        "success": True,
        "file": file_path.name,
        "chunks_created": embedded_count,
        "total_chunks": len(chunks)
    }


def ingest_training_corpus(clear_existing: bool = True) -> Dict[str, Any]:
    """
    Ingest all files from the Training folder into the knowledge RAG.

    Args:
        clear_existing: If True, clear existing knowledge data first

    Returns:
        Summary of ingestion results
    """
    if not TRAINING_DIR.exists():
        return {"success": False, "error": f"Training directory not found: {TRAINING_DIR}"}

    collection = get_collection(KNOWLEDGE_COLLECTION)
    if not collection:
        return {"success": False, "error": "Collection not available"}

    # Clear existing knowledge data if requested
    if clear_existing:
        try:
            # Get all IDs that start with "knowledge_"
            existing = collection.get(where={"source_type": "knowledge_corpus"})
            if existing and existing.get('ids'):
                collection.delete(ids=existing['ids'])
                logger.info(f"Cleared {len(existing['ids'])} existing knowledge chunks")
        except Exception as e:
            logger.warning(f"Could not clear existing data: {e}")

    # Find all supported files
    supported_extensions = {'.pdf', '.docx', '.pptx', '.xlsx'}
    files = [f for f in TRAINING_DIR.iterdir()
             if f.is_file() and f.suffix.lower() in supported_extensions]

    results = {
        "success": True,
        "files_processed": 0,
        "files_failed": 0,
        "total_chunks": 0,
        "details": []
    }

    for file_path in files:
        logger.info(f"Ingesting: {file_path.name}")
        result = ingest_file(file_path)

        if result.get("success"):
            results["files_processed"] += 1
            results["total_chunks"] += result.get("chunks_created", 0)
        else:
            results["files_failed"] += 1

        results["details"].append(result)

    return results


def get_corpus_stats() -> Dict[str, Any]:
    """Get statistics about the knowledge corpus."""
    collection = get_collection(KNOWLEDGE_COLLECTION)
    if not collection:
        return {"available": False}

    try:
        # Get all knowledge corpus entries
        existing = collection.get(where={"source_type": "knowledge_corpus"})

        # Count unique source files
        source_files = set()
        if existing and existing.get('metadatas'):
            for meta in existing['metadatas']:
                if meta.get('source_file'):
                    source_files.add(meta['source_file'])

        return {
            "available": True,
            "collection": KNOWLEDGE_COLLECTION,
            "total_chunks": len(existing.get('ids', [])) if existing else 0,
            "source_files": list(source_files),
            "file_count": len(source_files)
        }
    except Exception as e:
        logger.error(f"Error getting corpus stats: {e}")
        return {"available": False, "error": str(e)}

"""
Unified file parsing utilities.

This module consolidates parsing functionality from engine.py and corpus.py:
- Structured parsing (returns dict with headers, rows, metadata) for inventory processing
- Text extraction (returns plain text) for embedding/RAG

Supported formats:
- Excel (.xlsx, .xls)
- CSV (.csv)
- PDF (.pdf)
- Word (.docx)
- PowerPoint (.pptx)
"""
import csv
import hashlib
import logging
import re
import xml.etree.ElementTree as ET
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

import pdfplumber

logger = logging.getLogger(__name__)

# XML namespaces for Excel parsing
NS_MAIN = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
OFFICE_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# =============================================================================
# Utility Functions
# =============================================================================

def sha256_path(path: Path) -> str:
    """Calculate SHA256 hash of a file."""
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def date_from_filename_or_mtime(path: Path) -> datetime:
    """Extract date from filename pattern or fall back to mtime."""
    name = path.stem
    # Match patterns like "12.31.2024" or "1.5.24"
    m = re.search(r"(\d{1,2})\.(\d{1,2})\.(\d{2,4})", name)
    if m:
        mm, dd, yy = m.groups()
        if len(yy) == 2:
            yy = "20" + yy
        try:
            return datetime(int(yy), int(mm), int(dd))
        except ValueError:
            pass
    return datetime.fromtimestamp(path.stat().st_mtime)


def normalize_text(text: str) -> str:
    """Normalize text for comparison (lowercase, remove special chars)."""
    if not text:
        return ""
    t = text.lower().strip()
    t = re.sub(r"[^a-z0-9\s]+", " ", t)
    t = re.sub(r"\s+", " ", t)
    return t


def to_float(val: Any) -> Optional[float]:
    """Convert value to float, handling currency and comma separators."""
    if val == "" or val is None:
        return None
    if isinstance(val, (int, float)):
        return float(val)
    try:
        cleaned = str(val).replace(",", "").replace("$", "").strip()
        return float(cleaned)
    except (ValueError, TypeError):
        return None


# =============================================================================
# Excel Parsing (Low-level XML-based)
# =============================================================================

def _col_index(cell_ref: str) -> int:
    """Convert Excel column reference (A, B, AA, etc.) to 1-based index."""
    col = ""
    for ch in cell_ref:
        if ch.isalpha():
            col += ch
        else:
            break
    idx = 0
    for c in col:
        idx = idx * 26 + (ord(c.upper()) - 64)
    return idx


def _read_shared_strings(zf: zipfile.ZipFile) -> List[str]:
    """Read shared strings table from Excel file."""
    try:
        data = zf.read("xl/sharedStrings.xml")
    except KeyError:
        return []
    root = ET.fromstring(data)
    strings = []
    for si in root.findall(f"{{{NS_MAIN}}}si"):
        texts = [t.text or "" for t in si.findall(f".//{{{NS_MAIN}}}t")]
        strings.append("".join(texts))
    return strings


def _read_workbook_sheets(zf: zipfile.ZipFile) -> List[Tuple[str, str]]:
    """Read sheet names and paths from Excel workbook."""
    wb = ET.fromstring(zf.read("xl/workbook.xml"))
    rels = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))

    id_to_target = {}
    for rel in rels.findall(f"{{{REL_NS}}}Relationship"):
        id_to_target[rel.attrib["Id"]] = rel.attrib["Target"]

    sheets = []
    for sh in wb.findall(f"{{{NS_MAIN}}}sheets/{{{NS_MAIN}}}sheet"):
        name = sh.attrib.get("name")
        rid = sh.attrib.get(f"{{{REL_NS}}}id")
        if not rid:
            rid = sh.attrib.get(f"{{{OFFICE_REL_NS}}}id")

        target = id_to_target.get(rid, "")
        if target and not target.startswith("xl/"):
            target = "xl/" + target
        sheets.append((name, target))
    return sheets


def _get_cell_value(cell: ET.Element, shared_strings: List[str]) -> Any:
    """Extract cell value from Excel XML element."""
    ctype = cell.attrib.get("t")
    v = cell.find(f"{{{NS_MAIN}}}v")

    if v is None and ctype == "inlineStr":
        t = cell.find(f"{{{NS_MAIN}}}is/{{{NS_MAIN}}}t")
        val = (t.text or "") if t is not None else ""
    elif v is None:
        val = ""
    else:
        val = v.text or ""
        if ctype == "s":
            try:
                val = shared_strings[int(val)]
            except (ValueError, IndexError):
                pass

    if isinstance(val, str):
        val = val.strip()
    return val


def _parse_data_sheet(
    zf: zipfile.ZipFile,
    sheet_path: str,
    shared_strings: List[str]
) -> Optional[Tuple[Dict[int, str], List[Dict[int, Any]]]]:
    """Parse a single Excel sheet and return headers and data rows."""
    data = zf.read(sheet_path)
    root = ET.fromstring(data)
    sheet_data = root.find(f"{{{NS_MAIN}}}sheetData")

    if sheet_data is None:
        return None

    rows = list(sheet_data.findall(f"{{{NS_MAIN}}}row"))

    # Find header row (first row with mostly string values)
    header_row_idx = None
    for row in rows:
        r_idx = int(row.attrib.get("r", "0"))
        values = []
        str_count = 0
        num_count = 0

        for c in row.findall(f"{{{NS_MAIN}}}c"):
            val = _get_cell_value(c, shared_strings)
            if val != "":
                values.append(val)
                if isinstance(val, str) and not val.replace(".", "", 1).isdigit():
                    str_count += 1
                else:
                    num_count += 1

        if len(values) >= 3 and str_count >= num_count:
            header_row_idx = r_idx
            break

    if header_row_idx is None:
        return None

    # Extract headers and data rows
    headers = {}
    data_rows = []

    for row in rows:
        r_idx = int(row.attrib.get("r", "0"))
        row_cells = {}

        for c in row.findall(f"{{{NS_MAIN}}}c"):
            cref = c.attrib.get("r")
            if not cref:
                continue
            val = _get_cell_value(c, shared_strings)
            row_cells[_col_index(cref)] = val

        if r_idx == header_row_idx:
            for idx, val in row_cells.items():
                if isinstance(val, str) and val:
                    headers[idx] = val
        elif r_idx > header_row_idx:
            if row_cells:
                data_rows.append(row_cells)

    return headers, data_rows


# =============================================================================
# Structured Parsing (returns dict with headers, rows, metadata)
# =============================================================================

def parse_excel(file_path: Union[str, Path]) -> Dict[str, Any]:
    """
    Parse an Excel file and return structured data.

    Returns:
        dict with 'headers', 'rows', 'metadata'
    """
    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    if p.suffix.lower() not in ('.xlsx', '.xls'):
        raise ValueError(f"Unsupported file type: {p.suffix}")

    result = {
        "headers": [],
        "rows": [],
        "metadata": {
            "filename": p.name,
            "file_size": p.stat().st_size,
            "parsed_at": datetime.now().isoformat(),
            "file_type": "excel"
        }
    }

    try:
        with zipfile.ZipFile(p, "r") as zf:
            shared = _read_shared_strings(zf)
            sheets = _read_workbook_sheets(zf)

            # Try to find a data sheet
            data_sheet = None
            for name, target in sheets:
                if name and name.lower().startswith("data"):
                    data_sheet = (name, target)
                    break

            # Fall back to first sheet
            if not data_sheet and sheets:
                data_sheet = sheets[0]

            if not data_sheet:
                return result

            parsed = _parse_data_sheet(zf, data_sheet[1], shared)
            if not parsed:
                return result

            headers_dict, data_rows = parsed

            # Convert headers to list
            max_col = max(headers_dict.keys()) if headers_dict else 0
            headers = []
            for i in range(1, max_col + 1):
                headers.append(headers_dict.get(i, f"Column_{i}"))
            result["headers"] = headers

            # Convert rows to list of dicts
            for row_data in data_rows:
                row = {}
                for col_idx, value in row_data.items():
                    header = headers_dict.get(col_idx, f"Column_{col_idx}")
                    row[header] = value
                if row:
                    result["rows"].append(row)

            result["metadata"]["sheet_name"] = data_sheet[0]
            result["metadata"]["row_count"] = len(result["rows"])
            result["metadata"]["column_count"] = len(headers)

    except zipfile.BadZipFile:
        raise ValueError(f"Invalid Excel file: {file_path}")
    except Exception as e:
        raise ValueError(f"Failed to parse file: {str(e)}")

    return result


def parse_csv(file_path: Union[str, Path]) -> Dict[str, Any]:
    """
    Parse a CSV file and return structured data.

    Returns:
        dict with 'headers', 'rows', 'metadata'
    """
    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    result = {
        "headers": [],
        "rows": [],
        "metadata": {
            "filename": p.name,
            "file_size": p.stat().st_size,
            "parsed_at": datetime.now().isoformat(),
            "file_type": "csv"
        }
    }

    try:
        with open(p, 'r', newline='', encoding='utf-8-sig') as f:
            # Try to detect delimiter
            sample = f.read(4096)
            f.seek(0)

            try:
                dialect = csv.Sniffer().sniff(sample)
            except csv.Error:
                dialect = csv.excel

            reader = csv.DictReader(f, dialect=dialect)
            result["headers"] = reader.fieldnames or []

            for row in reader:
                result["rows"].append(dict(row))

            result["metadata"]["row_count"] = len(result["rows"])
            result["metadata"]["column_count"] = len(result["headers"])

    except Exception as e:
        raise ValueError(f"Failed to parse CSV: {str(e)}")

    return result


def parse_pdf(file_path: Union[str, Path]) -> Dict[str, Any]:
    """
    Parse a PDF file and return structured data.
    Extracts both text content and tables from all pages.

    Returns:
        dict with 'headers', 'rows', 'metadata', 'text_content'
    """
    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    if p.suffix.lower() != '.pdf':
        raise ValueError(f"Unsupported file type: {p.suffix}")

    result = {
        "headers": [],
        "rows": [],
        "text_content": [],
        "metadata": {
            "filename": p.name,
            "file_size": p.stat().st_size,
            "parsed_at": datetime.now().isoformat(),
            "file_type": "pdf"
        }
    }

    try:
        with pdfplumber.open(p) as pdf:
            result["metadata"]["page_count"] = len(pdf.pages)
            all_text = []

            for page_num, page in enumerate(pdf.pages, 1):
                # Extract text from the page
                text = page.extract_text()
                if text:
                    all_text.append({
                        "page": page_num,
                        "content": text.strip()
                    })

                # Extract tables from the page
                tables = page.extract_tables()
                for table_idx, table in enumerate(tables):
                    if not table or len(table) < 2:
                        continue

                    # First row is typically headers
                    headers = [
                        str(cell).strip() if cell else f"Column_{i}"
                        for i, cell in enumerate(table[0])
                    ]

                    # Track headers (use first table's headers as primary)
                    if not result["headers"]:
                        result["headers"] = headers

                    # Convert remaining rows to dicts
                    for row in table[1:]:
                        if not row or all(cell is None or str(cell).strip() == "" for cell in row):
                            continue
                        row_dict = {}
                        for i, cell in enumerate(row):
                            header = headers[i] if i < len(headers) else f"Column_{i}"
                            row_dict[header] = str(cell).strip() if cell else ""
                        result["rows"].append(row_dict)

            result["text_content"] = all_text
            result["metadata"]["row_count"] = len(result["rows"])
            result["metadata"]["column_count"] = len(result["headers"])
            result["metadata"]["text_blocks"] = len(all_text)

    except Exception as e:
        raise ValueError(f"Failed to parse PDF: {str(e)}")

    return result


def parse_file(file_path: Union[str, Path]) -> Dict[str, Any]:
    """
    Universal file parser - routes to appropriate parser based on file type.
    Returns structured data (headers, rows, metadata).

    For text-only extraction, use extract_text() instead.

    Returns:
        dict with 'headers', 'rows', 'metadata' (and 'text_content' for PDFs)
    """
    p = Path(file_path)
    suffix = p.suffix.lower()

    if suffix in ('.xlsx', '.xls'):
        return parse_excel(file_path)
    elif suffix == '.csv':
        return parse_csv(file_path)
    elif suffix == '.pdf':
        return parse_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


# =============================================================================
# Text Extraction (returns plain text for embedding/RAG)
# =============================================================================

def extract_text_from_pdf(file_path: Union[str, Path]) -> str:
    """Extract text from PDF file."""
    p = Path(file_path)
    text_parts = []

    try:
        with pdfplumber.open(p) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text.strip())
    except Exception as e:
        logger.error(f"Error parsing PDF {file_path}: {e}")

    return "\n\n".join(text_parts)


def extract_text_from_docx(file_path: Union[str, Path]) -> str:
    """Extract text from DOCX file."""
    from docx import Document as DocxDocument

    p = Path(file_path)
    text_parts = []

    try:
        doc = DocxDocument(p)
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        # Also get text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    text_parts.append(row_text)
    except Exception as e:
        logger.error(f"Error parsing DOCX {file_path}: {e}")

    return "\n\n".join(text_parts)


def extract_text_from_pptx(file_path: Union[str, Path]) -> str:
    """Extract text from PPTX file."""
    from pptx import Presentation

    p = Path(file_path)
    text_parts = []

    try:
        prs = Presentation(p)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    except Exception as e:
        logger.error(f"Error parsing PPTX {file_path}: {e}")

    return "\n\n".join(text_parts)


def extract_text_from_excel(file_path: Union[str, Path]) -> str:
    """Extract text from Excel file using openpyxl."""
    import openpyxl

    p = Path(file_path)
    text_parts = []

    try:
        wb = openpyxl.load_workbook(p, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[Sheet: {sheet_name}]"]
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    sheet_text.append(row_text)
            if len(sheet_text) > 1:
                text_parts.append("\n".join(sheet_text))
        wb.close()
    except Exception as e:
        logger.error(f"Error parsing XLSX {file_path}: {e}")

    return "\n\n".join(text_parts)


def extract_text(file_path: Union[str, Path]) -> Optional[str]:
    """
    Extract text content from any supported file type.
    Use this for embedding/RAG workflows.

    Supported formats: .pdf, .docx, .pptx, .xlsx

    Returns:
        Plain text content, or None if unsupported/failed
    """
    p = Path(file_path)
    suffix = p.suffix.lower()

    if suffix == '.pdf':
        return extract_text_from_pdf(p)
    elif suffix == '.docx':
        return extract_text_from_docx(p)
    elif suffix == '.pptx':
        return extract_text_from_pptx(p)
    elif suffix == '.xlsx':
        return extract_text_from_excel(p)
    else:
        logger.warning(f"Unsupported file type: {suffix} for {file_path}")
        return None


# =============================================================================
# Excel-specific utilities for site extraction
# =============================================================================

def extract_site_from_excel(
    file_path: Union[str, Path],
    slugify_func=None
) -> Optional[str]:
    """
    Extract site name from Excel file header rows.

    Looks for patterns like:
        "Site Name - Location (123) (VENDOR)" -> "site_name_location"
        "Building Name, Floor 100 (VENDOR)" -> "building_name_floor_100"

    Args:
        file_path: Path to Excel file
        slugify_func: Optional function to convert site name to ID

    Returns:
        Site ID string or None
    """
    if slugify_func is None:
        from .naming import slugify
        slugify_func = slugify

    p = Path(file_path)

    try:
        with zipfile.ZipFile(p, "r") as zf:
            shared = _read_shared_strings(zf)
            sheets = _read_workbook_sheets(zf)

            for name, target in sheets:
                site_id = _extract_site_from_sheet(zf, target, shared, slugify_func)
                if site_id:
                    return site_id
    except Exception:
        pass

    return None


def _extract_site_from_sheet(
    zf: zipfile.ZipFile,
    sheet_path: str,
    shared_strings: List[str],
    slugify_func
) -> Optional[str]:
    """Extract site name from early rows of a sheet."""
    try:
        data = zf.read(sheet_path)
        root = ET.fromstring(data)
        sheet_data = root.find(f"{{{NS_MAIN}}}sheetData")

        if sheet_data is None:
            return None

        # Check first 5 rows for site name
        rows = sheet_data.findall(f"{{{NS_MAIN}}}row")[:5]
        for row in rows:
            for c in row.findall(f"{{{NS_MAIN}}}c"):
                val = _get_cell_value(c, shared_strings)
                if not isinstance(val, str) or not val:
                    continue

                # Skip generic headers
                val_lower = val.lower()
                skip_terms = [
                    "inventory", "report", "property", "proprietary",
                    "current", "preferred", "printed by"
                ]
                if any(term in val_lower for term in skip_terms):
                    continue

                # Match pattern with (COMPASS) at end
                match = re.match(r'^(.+?)\s*\(COMPASS\)\s*$', val, re.IGNORECASE)
                if match:
                    site_name = match.group(1).strip()
                    site_name = re.sub(r'\s*\(\d+\)\s*$', '', site_name)
                    site_id = slugify_func(site_name)
                    if site_id and len(site_id) >= 2:
                        return site_id

                # Match pattern with (number)
                match = re.match(r'^([A-Za-z0-9\s\-,]+)\s*\(\d+\)', val)
                if match:
                    site_name = match.group(1).strip()
                    site_id = slugify_func(site_name)
                    if site_id and len(site_id) >= 2:
                        return site_id
    except Exception:
        pass

    return None
